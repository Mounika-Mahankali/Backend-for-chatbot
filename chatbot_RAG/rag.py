from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_core.language_models import LLM
from typing import Optional, List
from pptx import Presentation
from langchain_core.documents import Document
from multimodal import get_response
import pickle
import os
import re
from logger import log_execution

# Document loaders
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,
    CSVLoader
)

from langchain_text_splitters import RecursiveCharacterTextSplitter
from pathlib import Path
from rank_bm25 import BM25Okapi
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

vectorless_docs = {}   # session_id to docs
bm25_models = {}       # session_id to  bm25 model
CHUNKS_FILE = "vectorless_chunks.pkl"
# Load chunks if file exists or Load saved chunks
if os.path.exists(CHUNKS_FILE):
    with open(CHUNKS_FILE, "rb") as f:
        data = pickle.load(f)

        # correct structure (dictionary)
        if isinstance(data, dict):
            vectorless_docs = data.get("docs", {})
            bm25_models = data.get("bm25_models", {})
        else:
            
            print("Old chunk format detected. Resetting...")
            vectorless_docs = {}
            bm25_models = {}

        print("Loaded vectorless docs:", len(vectorless_docs))

# Custom LLM
class CustomLLM(LLM):

    def _call(self, prompt: str, stop: Optional[List[str]] = None) -> str:
        return get_response(prompt)

    @property
    def _identifying_params(self):
        return {}

    @property
    def _llm_type(self):
        return "custom_groq"


@log_execution("llm_generation")
def generate_response(llm, prompt):
    return llm.invoke(prompt)



@log_execution("reranking")
def rerank_documents(query, docs, top_k=15):

    if not docs:
        return docs

    # query_tokens = re.findall(r"\w+", query.lower())
    query_tokens = re.findall(r"\w+", query.lower())

    # add simple normalization
    expanded_tokens = set(query_tokens)

    # basic normalization (generic)
    for token in query_tokens:
        if token.endswith("s"):
            expanded_tokens.add(token[:-1])
        if token.endswith("ing"):
            expanded_tokens.add(token[:-3])

    query_tokens = list(expanded_tokens)
        

    scored_docs = []

    for doc in docs:

        text = doc.page_content.lower()

        score = sum(1 for token in query_tokens if token in text)

        if "\n" in doc.page_content:
            score += 2

       
        if any(char.isdigit() for char in doc.page_content):
            score += 2

          
        scored_docs.append((score, doc))

    ranked_docs = sorted(
        scored_docs,
        key=lambda x: x[0],
        reverse=True
    )

    return [doc for _, doc in ranked_docs[:top_k]]


@log_execution("image_llm_generation")
def generate_image_response(prompt, image):
    return get_response(prompt, image)


@log_execution("post_processing")
def post_process_response(response):
    return response.strip()


# Load documents
def load_pptx_documents(file_path):

    presentation = Presentation(str(file_path))
    documents = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_parts.append(shape.text.strip())

        slide_text = "\n".join(slide_parts).strip()

        if slide_text:
            documents.append(
                Document(
                    page_content=f"Slide {slide_number}\n{slide_text}",
                    metadata={
                        "slide_number": slide_number,
                        "content_type": "slide_text"
                    }
                )
            )

    return documents


def load_documents(file_path):

    file_path = Path(file_path)
    file_suffix = file_path.suffix.lower()
    file_str = str(file_path)

    if file_suffix == ".pdf":
        loader = PyPDFLoader(file_str)
        return loader.load()

    if file_suffix == ".docx":
        loader = Docx2txtLoader(file_str)
        return loader.load()

    if file_suffix == ".txt":
        loader = TextLoader(file_str)
        return loader.load()

    if file_suffix == ".csv":
        loader = CSVLoader(file_str)
        return loader.load()

    if file_suffix == ".pptx":
        return load_pptx(file_path)

    raise ValueError(f"Unsupported file format: {file_path}")


def extract_images_from_ppt(ppt_path):

    images = []

    prs = Presentation(str(ppt_path))

    for slide_number, slide in enumerate(prs.slides, start=1):
        image_index = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_index += 1

                image = shape.image
                image_bytes = image.blob

                images.append(
                    {
                        "bytes": image_bytes,
                        "slide_number": slide_number,
                        "image_index": image_index
                    }
                )

    return images


def persist_vectorless_store():

    with open(CHUNKS_FILE, "wb") as f:
        pickle.dump(
            {
                "docs": vectorless_docs,
                "bm25_models": bm25_models
            },
            f
        )


def rebuild_bm25_for_session(chat_session_id):

    docs = vectorless_docs.get(chat_session_id, [])

    if not docs:
        bm25_models.pop(chat_session_id, None)
        persist_vectorless_store()
        return

    tokenized_docs = [
        doc.page_content.lower().split()
        for doc in docs
    ]

    bm25_models[chat_session_id] = BM25Okapi(tokenized_docs)
    persist_vectorless_store()

#chunking 
def store_documents_in_session(documents, chat_session_id):

    if not documents:
        return

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1500,
    chunk_overlap=200
    )

    new_docs = text_splitter.split_documents(documents)

    if chat_session_id not in vectorless_docs:
        vectorless_docs[chat_session_id] = []

    vectorless_docs[chat_session_id].extend(new_docs)
    rebuild_bm25_for_session(chat_session_id)

# Load vectorless docs
def load_vectorless_docs(pdf_paths, chat_session_id):

    global vectorless_docs
    global bm25_model

    all_documents = []

    for file_path in pdf_paths:

        file_path = Path(file_path)

        documents = load_documents(file_path)

        if documents:
            for doc in documents:
                doc.metadata = {
                    **doc.metadata,
                    "source": file_path.name
                }
                all_documents.append(doc)

        # Extract images from PPT
        if file_path.suffix == ".pptx":

            images = extract_images_from_ppt(file_path)

            if images:
                print(f"Processing images...")

                for i, img in enumerate(images):
                    try:
                        caption = get_response(
                            "Describe this image in detail",
                            img["bytes"]
                        )

                        all_documents.append(
                            Document(
                                page_content=caption,
                                metadata={
                                    "source": file_path.name,
                                    "slide_number": img["slide_number"],
                                    "image_index": img["image_index"],
                                    "content_type": "slide_image_caption"
                                }
                            )
                        )

                        print(f"Processed {i+1}/{len(images)} images")

                    except Exception as e:
                        print(f"Image extraction error: {e}")

    store_documents_in_session(all_documents, chat_session_id)

    print("Chunks stored:", len(vectorless_docs))


def save_image_description_to_db(image_bytes, filename, chat_session_id):

    caption = generate_image_response("Describe this image in detail", image_bytes)

    document = Document(
        page_content=caption,
        metadata={
            "source": filename,
            "content_type": "uploaded_image_caption"
        }
    )

    store_documents_in_session([document], chat_session_id)

    return caption


def filter_best_document(docs):

    if not docs:
        return docs

    doc_scores = {}

    for doc in docs:
        source = doc.metadata.get("source", "Unknown")

        if source not in doc_scores:
            doc_scores[source] = 0

        doc_scores[source] += 1

    if not doc_scores:
        return docs

    best_source = max(doc_scores, key=doc_scores.get)

    filtered_docs = [
        doc for doc in docs
        if doc.metadata.get("source") == best_source
    ]

    return filtered_docs


# Load RAG
def load_rag(chat_session_id=None):

    print("Loading Vector-less RAG...")

    llm = CustomLLM()

    def format_docs(docs):
        return "\n\n".join(doc.page_content for doc in docs)

    def vectorless_rag(query):

        docs = hybrid_retrieval(query, chat_session_id)

        query_words = set(query.lower().split())

        filtered_docs = []
        for doc in docs:
            text = doc.page_content.lower()
            
            match_count = sum(1 for word in query_words if word in text)
            
            if match_count >= 1:
                filtered_docs.append(doc)

        # fallback if nothing found
        if filtered_docs:
            docs = filtered_docs

        
        if len(query.split()) <= 4:
            docs = docs[:20]   

        docs = rerank_documents(query, docs)


        #docs = filter_best_document(docs)

        if not docs:
            response = generate_response(llm, query)
            return post_process_response(response)
        
        docs = docs[:12]
        context = "\n\n".join(doc.page_content for doc in docs[:8])

        
        prompt = create_rag_prompt(context, query)

        response = generate_response(llm, prompt)

        response = post_process_response(response)

        best_source = docs[0].metadata.get("source", "Unknown")

        response += "\n\n**Sources:**\n"
        response += f"- {best_source}\n"

        return response
    return vectorless_rag


@log_execution("prompt_creation")
def create_prompt(context, query):
    return f"""
You are a helpful AI assistant.

Use the provided context to answer the question completely.
If the answer exists in context, return full information.

If context is not relevant, answer using general knowledge.

Context:
{context}

Question:
{query}

Answer:
"""


# Save docs when uploaded
def save_pdfs_to_db(pdf_paths, chat_session_id):

    load_vectorless_docs(pdf_paths, chat_session_id)

def get_persistent_retriever(chat_session_id):
    return None

@log_execution("bm25_retrieval")
def bm25_retrieval(query, chat_session_id):

    global vectorless_docs
    global bm25_model

    if chat_session_id not in vectorless_docs:
        return []

    query_tokens = preprocess_query(query)
    bm25_model = bm25_models.get(chat_session_id)

    if not bm25_model:
        return []

    scores = bm25_model.get_scores(query_tokens)
    scores = list(scores)

    scored_docs = list(zip(scores, vectorless_docs[chat_session_id]))

    if not scored_docs:
        return []

    scored_docs = sorted(
        scored_docs,
        key=lambda x: x[0],
        reverse=True
    )

    
    top_score = scored_docs[0][0] if scored_docs else 0

    # filtered_docs = [
    #     doc for score, doc in scored_docs
    #     if score >= top_score * 0.6   
    # ]
    filtered_docs = [doc for score, doc in scored_docs]

    return filtered_docs[:15]


def preprocess_query(query):

    stopwords = {
        "what","is","are","the","of","in","for",
        "does","do","a","an","and","to"
    }

    tokens = re.findall(r"\w+", query.lower())

    filtered = [
        token for token in tokens
        if token not in stopwords
    ]

    return filtered




def keyword_retrieval(query, chat_session_id):

    global vectorless_docs

    query_tokens = re.findall(r"\w+", query.lower())

    matched_docs = []

    docs = vectorless_docs.get(chat_session_id, [])

    for doc in docs:
        text = doc.page_content.lower()

        if any(token in text for token in query_tokens):
            matched_docs.append(doc)

    return matched_docs[:15]


def hybrid_retrieval(query, chat_session_id):

    top_k = dynamic_top_k(query)

    bm25_docs = bm25_retrieval(query, chat_session_id)

    keyword_docs = keyword_retrieval(query, chat_session_id)

    combined_docs = bm25_docs + keyword_docs

    unique_docs = []
    seen = set()

    for doc in combined_docs:
        content = doc.page_content

        if content not in seen:
            unique_docs.append(doc)
            seen.add(content)

    return unique_docs[:15]


#diff queries need different retrieval sizes 
# def dynamic_top_k(query):

#     query = query.lower()

#     if any(word in query for word in ["list", "keywords", "types", "advantages"]):
#         return 6

#     if any(word in query for word in ["what is", "define", "meaning"]):
#         return 2

#     if any(word in query for word in ["difference", "compare"]):
#         return 4

#     return 5

def dynamic_top_k(query):
    return 12
def create_rag_prompt(context, query):
    return f"""
You are an AI assistant.

Extract the exact information from the context.

If the answer is present as a table or structured format:
- return it clearly
- do NOT summarize
- do NOT say "not available" if present

If multiple rows exist, include all rows.

Context:
{context}

Question:
{query}

Answer:
"""


def load_pptx(file_path):
    docs = []

    prs = Presentation(file_path)

    for slide in prs.slides:
        slide_text = []

        for shape in slide.shapes:

            # Normal text
            if hasattr(shape, "text"):
                slide_text.append(shape.text)

            # Table extraction (VERY IMPORTANT)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_data = " | ".join(cell.text.strip() for cell in row.cells)
                    slide_text.append(row_data)

        full_text = "\n".join(slide_text)

        docs.append(Document(page_content=full_text))

    return docs

